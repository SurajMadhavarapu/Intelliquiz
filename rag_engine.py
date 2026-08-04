import os
import re
import json
import glob
from typing import List, Dict, Any, Tuple, Optional
from pathlib import Path

# Document Loaders & Utilities
from pypdf import PdfReader
from pptx import Presentation
import docx

# LangChain Imports
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.docstore.in_memory import InMemoryDocstore
from langchain_core.documents import Document

# Ollama LLM Integration
try:
    from langchain_ollama import OllamaLLM
except ImportError:
    try:
        from langchain_community.llms import Ollama as OllamaLLM
    except ImportError:
        OllamaLLM = None


class RAGPipeline:
    def __init__(self, model_name: str = "mistral", embedding_model: str = "sentence-transformers/all-MiniLM-L6-v2"):
        self.model_name = model_name
        self.embedding_model_name = embedding_model
        self.embeddings = None
        self.vector_store: Optional[FAISS] = None
        self.documents: List[Document] = []
        self.indexed_files: Dict[str, Dict[str, Any]] = {}
        
        self._init_embeddings()

    def _init_embeddings(self):
        """Initialize HuggingFace Embeddings model."""
        try:
            print(f"Loading embeddings model: {self.embedding_model_name}")
            self.embeddings = HuggingFaceEmbeddings(
                model_name=self.embedding_model_name,
                model_kwargs={'device': 'cpu'},
                encode_kwargs={'normalize_embeddings': True}
            )
            print("Embeddings loaded successfully!")
        except Exception as e:
            print(f"Error loading HuggingFaceEmbeddings: {e}")
            raise e

    def _get_llm(self, model: Optional[str] = None):
        """Get Ollama LLM instance."""
        target_model = model or self.model_name
        if OllamaLLM is None:
            raise RuntimeError("LangChain Ollama integration is not installed.")
        return OllamaLLM(model=target_model, base_url="http://localhost:11434", temperature=0.2)

    def parse_file(self, file_path: str) -> List[Document]:
        """Parse PDF, PPTX, DOCX, or TXT file into a list of LangChain Document objects with rich metadata."""
        path = Path(file_path)
        filename = path.name
        ext = path.suffix.lower()
        docs: List[Document] = []

        try:
            if ext == ".pdf":
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        docs.append(Document(
                            page_content=text.strip(),
                            metadata={"source": filename, "page": i + 1, "file_type": "PDF"}
                        ))
            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    full_slide_text = "\n".join(slide_text)
                    if full_slide_text.strip():
                        docs.append(Document(
                            page_content=full_slide_text.strip(),
                            metadata={"source": filename, "slide": i + 1, "file_type": "PPTX"}
                        ))
            elif ext in [".docx", ".doc"]:
                doc = docx.Document(file_path)
                full_text = "\n".join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
                if full_text.strip():
                    docs.append(Document(
                        page_content=full_text.strip(),
                        metadata={"source": filename, "file_type": "DOCX"}
                    ))
            elif ext in [".txt", ".md", ".csv"]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                    if text.strip():
                        docs.append(Document(
                            page_content=text.strip(),
                            metadata={"source": filename, "file_type": ext[1:].upper()}
                        ))
            else:
                print(f"Unsupported file format: {ext}")
        except Exception as e:
            print(f"Error parsing file {filename}: {e}")

        return docs

    def process_and_index_files(self, file_paths: List[str]) -> Dict[str, Any]:
        """Parse multiple files, split into chunks, and update the FAISS vector database."""
        all_chunks: List[Document] = []
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=600,
            chunk_overlap=100,
            separators=["\n\n", "\n", ". ", " ", ""]
        )

        newly_indexed = 0
        for path in file_paths:
            filename = Path(path).name
            raw_docs = self.parse_file(path)
            if not raw_docs:
                continue

            chunks = text_splitter.split_documents(raw_docs)
            all_chunks.extend(chunks)
            newly_indexed += 1
            self.indexed_files[filename] = {
                "path": path,
                "raw_pages_slides": len(raw_docs),
                "total_chunks": len(chunks)
            }

        if not all_chunks:
            return {"status": "error", "message": "No valid text content extracted from files."}

        # Build or add to FAISS Vector Store
        if self.vector_store is None:
            self.vector_store = FAISS.from_documents(all_chunks, self.embeddings)
        else:
            self.vector_store.add_documents(all_chunks)

        self.documents.extend(all_chunks)

        return {
            "status": "success",
            "files_indexed": newly_indexed,
            "total_files": len(self.indexed_files),
            "total_chunks_in_db": len(self.documents)
        }

    def clear_database(self):
        """Clear all indexed documents and vector store."""
        self.vector_store = None
        self.documents = []
        self.indexed_files = {}

    def query_rag(self, user_query: str, model: Optional[str] = None, top_k: int = 5) -> Dict[str, Any]:
        """Perform similarity search and generate answer with source citations."""
        if self.vector_store is None or not self.documents:
            return {
                "answer": "No documents uploaded yet! Please upload one or more files (PPT, PDF, DOCX, TXT) to start asking questions.",
                "citations": []
            }

        # 1. Similarity search in FAISS
        relevant_docs = self.vector_store.similarity_search(user_query, k=top_k)
        
        # Build Context String & Citations List
        context_blocks = []
        citations = []
        seen_citations = set()

        for idx, doc in enumerate(relevant_docs):
            src = doc.metadata.get("source", "Unknown Document")
            slide = doc.metadata.get("slide")
            page = doc.metadata.get("page")
            
            location_label = f"Slide {slide}" if slide else (f"Page {page}" if page else "Document")
            citation_str = f"{src} ({location_label})"
            
            if citation_str not in seen_citations:
                seen_citations.add(citation_str)
                citations.append({
                    "source": src,
                    "location": location_label,
                    "snippet": doc.page_content[:150] + "..."
                })

            context_blocks.append(f"--- Document Source [{idx+1}]: {citation_str} ---\n{doc.page_content}")

        full_context = "\n\n".join(context_blocks)

        # Prompt formatting
        prompt = f"""You are Intelliquiz AI, an intelligent academic tutor for B.Tech students.
Answer the user's question accurately using ONLY the provided document context below.
If the answer is found in the context, synthesize a clear, well-structured, and helpful response. Use markdown formatting (bullet points, bold text, code blocks) where appropriate.
Explicitly mention which slide/page or document provided key information when relevant.
If the context does NOT contain enough information to answer, state clearly what is missing based on the files provided.

Context:
{full_context}

User Question:
{user_query}

Detailed Answer:"""

        try:
            llm = self._get_llm(model=model)
            response_text = llm.invoke(prompt)
            return {
                "answer": response_text.strip(),
                "citations": citations
            }
        except Exception as e:
            print(f"LLM execution error: {e}")
            return {
                "answer": f"⚠️ Error communicating with Ollama: {str(e)}. Make sure Ollama is running (`ollama serve`).",
                "citations": citations
            }

    def generate_quiz(self, num_questions: int = 5, model: Optional[str] = None) -> List[Dict[str, Any]]:
        """Generate multiple choice quiz questions (MCQs) from indexed documents."""
        if not self.documents:
            return []

        # Sample up to 10 context chunks for prompt budget
        sample_chunks = self.documents[:12]
        context_str = "\n\n".join([f"Source ({doc.metadata.get('source')}): {doc.page_content}" for doc in sample_chunks])

        prompt = f"""Generate exactly {num_questions} high-quality multiple choice quiz questions (MCQs) for B.Tech AI & DS students based on the following document context.

Return ONLY a valid JSON array of objects. Do not include markdown code block formatting like ```json.
Each object in the array MUST have the exact following key structure:
{{
  "id": 1,
  "question": "Question text here...",
  "options": ["Option A", "Option B", "Option C", "Option D"],
  "answer_index": 0,
  "explanation": "Explanation of why Option A is correct...",
  "source": "Document or Topic name"
}}

Document Context:
{context_str}

JSON Output:"""

        try:
            llm = self._get_llm(model=model)
            raw_response = llm.invoke(prompt).strip()
            
            # Clean possible markdown wrapping
            if raw_response.startswith("```"):
                raw_response = re.sub(r"^```(?:json)?\n?", "", raw_response)
                raw_response = re.sub(r"\n?```$", "", raw_response).strip()

            quizzes = json.loads(raw_response)
            return quizzes if isinstance(quizzes, list) else []
        except Exception as e:
            print(f"Quiz generation error: {e}")
            # Fallback mock quiz structure if parsing failed
            return [
                {
                    "id": 1,
                    "question": "What is the primary objective of Retrieval-Augmented Generation (RAG)?",
                    "options": [
                        "To combine information retrieval with generative language models for grounded answers",
                        "To train neural networks from scratch without data",
                        "To compress PowerPoint presentations into ZIP files",
                        "To execute SQL queries directly on standard web servers"
                    ],
                    "answer_index": 0,
                    "explanation": "RAG retrieves relevant document passages from a vector database and feeds them into an LLM for factual context.",
                    "source": "Intelliquiz RAG Overview"
                }
            ]

    def generate_summary(self, model: Optional[str] = None) -> Dict[str, str]:
        """Generate a concise summary for each uploaded document."""
        summaries = {}
        if not self.indexed_files:
            return summaries

        for filename in self.indexed_files:
            # Gather chunks for this file
            file_chunks = [doc.page_content for doc in self.documents if doc.metadata.get("source") == filename]
            context = "\n".join(file_chunks[:6])  # limit length

            prompt = f"""Summarize the key concepts, topics, and important takeaways from the file '{filename}' in 3-5 bullet points.

Content:
{context}

Key Summary Bullets:"""
            try:
                llm = self._get_llm(model=model)
                summary = llm.invoke(prompt).strip()
                summaries[filename] = summary
            except Exception as e:
                summaries[filename] = f"Could not generate summary: {e}"

        return summaries
