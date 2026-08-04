from pptx import Presentation
from pptx.util import Inches, Pt
import os

prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Deep Learning & Neural Networks"
subtitle1.text = "B.Tech AI & DS Coursework - Slide Set 1"

# Slide 2: Convolutional Neural Networks
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
body2 = slide2.placeholders[1]
title2.text = "Convolutional Neural Networks (CNNs)"
body2.text = (
    "• CNNs are specialized neural architectures for grid-like data (images, spatial signals).\n"
    "• Key Layers:\n"
    "  1. Convolutional Layer: Applies kernels/filters for feature extraction (edges, textures).\n"
    "  2. Pooling Layer (Max Pooling / Average Pooling): Downsamples spatial dimensions.\n"
    "  3. Fully Connected (FC) Layer: Performs final classification."
)

# Slide 3: Transformer Architecture & Attention
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
body3 = slide3.placeholders[1]
title3.text = "Self-Attention and Transformers"
body3.text = (
    "• Proposed by Vaswani et al. (2017) in 'Attention Is All You Need'.\n"
    "• Key Mechanism: Scaled Dot-Product Attention.\n"
    "  Attention(Q, K, V) = softmax( (Q * K^T) / sqrt(d_k) ) * V\n"
    "• Multi-Head Attention allows joint projection to different representation subspaces."
)

os.makedirs("sample_data", exist_ok=True)
prs.save("sample_data/deep_learning_lecture.pptx")
print("Saved sample_data/deep_learning_lecture.pptx successfully!")
